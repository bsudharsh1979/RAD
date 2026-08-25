from __future__ import annotations

import hashlib
import json
import re
import uuid
from html import unescape
from pathlib import Path
from typing import Any

import nbformat

from app.config import settings
from app.db.models import (
    Course,
    EvidenceType,
    LearningObjective,
    Notebook,
    NotebookCell,
    SourceArtifact,
    SourceSpan,
)
from app.db.session import SessionLocal
from app.domains.retrieval.embeddings import hash_embedding
from app.ingestion.safety import inspect_code

HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)")
HTML_RE = re.compile(r"<[^>]+>")
OBJ_RE = re.compile(r"Learning Objectives:.*?(?=\n#{1,3}\s|\n-----|\Z)", re.S | re.I)


def _strip(md: str) -> str:
    return unescape(HTML_RE.sub("", md or "")).strip()


def _sha(path: Path) -> str:
    h = hashlib.sha256()
    h.update(path.read_bytes())
    return h.hexdigest()


def _id(prefix: str, *parts: Any) -> str:
    raw = "|".join(str(p) for p in parts)
    return prefix + hashlib.md5(raw.encode()).hexdigest()[:16]


def _cell_text(cell: Any) -> str:
    src = cell.get("source", "")
    if isinstance(src, list):
        return "".join(src)
    return str(src or "")


def _outputs(cell: Any) -> str:
    chunks: list[str] = []
    for o in cell.get("outputs") or []:
        ot = o.get("output_type")
        if ot == "stream":
            t = o.get("text", "")
            chunks.append("".join(t) if isinstance(t, list) else str(t))
        elif ot in ("execute_result", "display_data"):
            data = o.get("data") or {}
            t = data.get("text/plain", "")
            chunks.append("".join(t) if isinstance(t, list) else str(t))
        elif ot == "error":
            chunks.append(f"{o.get('ename')}: {o.get('evalue')}")
    return "\n".join(chunks).strip()


def _heading(text: str) -> str | None:
    for line in text.splitlines():
        m = HEADING_RE.match(_strip(line))
        if m:
            return HTML_RE.sub("", m.group(2)).strip()
    return None


NOTEBOOK_META: dict[str, dict[str, str]] = {
    "01_llm_intro.ipynb": {
        "purpose": "Load a HuggingFace fill-mask pipeline and peel it into tokenizer + model.",
        "why": "Every later notebook assumes you can treat transformers as preprocess → forward → postprocess.",
        "outcome": "Explain why pipelines hide tensors and why GPU RAM still matters.",
    },
    "02_llm_intake.ipynb": {
        "purpose": "Inspect BERT embeddings and self-attention.",
        "why": "Token IDs, 768-D embeddings, and Q/K/V are the intake contract for encoder models.",
        "outcome": "Describe word + position + type addition and multi-head self-attention.",
    },
    "03_encoder_task.ipynb": {
        "purpose": "Attach task heads: MLM, span QA, sequence classification, zero-shot.",
        "why": "The encoder body is reused; the head and number of queries change the task.",
        "outcome": "Choose token vs span vs sequence vs multi-query formulations.",
    },
    "04_seq2seq.ipynb": {
        "purpose": "Use T5/Flan-T5 for unbounded generation with cross-attention.",
        "why": "Encoders cannot emit novel tokens; decoders generate one token at a time.",
        "outcome": "Explain encoder-once / decoder-many and in-context learning vs T5 memorized prefixes.",
    },
    "05_multimodal.ipynb": {
        "purpose": "Treat audio and images as sequences that condition a text decoder.",
        "why": "Cross-attention is modality-agnostic if the encoder emits a sequence.",
        "outcome": "Contrast Whisper, ViT-GPT2/BLIP captioning, and CLIP dual encoders.",
    },
    "06_textgen.ipynb": {
        "purpose": "Move to decoder-only GPT/CodeGen/Llama-2 and quantization.",
        "why": "Chat models need templates, and 70B models do not fit consumer GPUs in FP16.",
        "outcome": "Use Llama-2 INST/SYS format and distinguish GPTQ vs bitsandbytes tradeoffs.",
    },
    "07_stateful_models.ipynb": {
        "purpose": "Wrap LLMs in LangChain memory, RAG, and ReAct agents.",
        "why": "A chat LLM without state is a stochastic parrot of the current prompt.",
        "outcome": "Separate conversation buffers, retrieval, and agent event loops.",
    },
    "08_assessment.ipynb": {
        "purpose": "Build a custom LangChain agent that talks to the user as a tool.",
        "why": "Certification requires composing course pieces, not recalling definitions.",
        "outcome": "Implement at least 3 of 5: memory, image, code, toxicity, emotion.",
    },
}


def ingest_course_materials(db=None) -> dict[str, Any]:
    own = db is None
    if own:
        db = SessionLocal()
    try:
        return _ingest(db)
    finally:
        if own:
            db.close()


def _ingest(db) -> dict[str, Any]:
    course = db.get(Course, "rad-llm")
    if not course:
        course = Course(
            id="rad-llm",
            slug="nvidia-dli-rad-llm",
            title=settings.course_title,
            description=(
                "Hands-on HuggingFace, encoder/decoder architectures, multimodal "
                "transformers, quantized Llama-2, LangChain orchestration, and a custom agent assessment."
            ),
        )
        db.add(course)
        db.commit()
        db.refresh(course)

    root: Path = settings.course_materials_dir
    stats = {"notebooks": 0, "spans": 0, "pdfs": 0, "pptx": 0, "skipped_exec": 0}
    if not root.exists():
        db.commit()
        return stats

    notebooks = sorted(root.rglob("*.ipynb"))
    for order, path in enumerate(notebooks, start=1):
        _ingest_notebook(db, path, order)
        stats["notebooks"] += 1

    for path in sorted(root.rglob("*.pdf")):
        _ingest_pdf(db, path)
        stats["pdfs"] += 1
    for path in sorted(list(root.rglob("*.pptx")) + list(root.rglob("*.ppt"))):
        _ingest_pptx(db, path)
        stats["pptx"] += 1

    stats["spans"] = db.query(SourceSpan).count()
    db.commit()
    return stats


def _ingest_notebook(db, path: Path, order: int) -> None:
    rel = str(path.relative_to(settings.course_materials_dir))
    sha = _sha(path)
    art_id = _id("art", rel)
    existing = db.get(SourceArtifact, art_id)
    if existing and existing.sha256 == sha:
        return

    nb = nbformat.read(path.open(), as_version=4)
    title = path.name
    lo_text = ""
    for cell in nb.cells:
        if cell.cell_type != "markdown":
            continue
        text = _strip(_cell_text(cell))
        if not title or title == path.name:
            h = _heading(_cell_text(cell))
            if h:
                title = h
        m = OBJ_RE.search(_cell_text(cell))
        if m:
            lo_text = _strip(m.group(0))

    if existing:
        nb_id = _id("nb", rel)
        db.query(NotebookCell).filter(NotebookCell.notebook_id == nb_id).delete()
        db.query(Notebook).filter(Notebook.id == nb_id).delete()
        db.query(SourceSpan).filter(SourceSpan.artifact_id == art_id).delete()
        db.delete(existing)
        db.flush()

    artifact = SourceArtifact(
        id=art_id,
        course_id="rad-llm",
        source_type="notebook",
        filename=path.name,
        relpath=rel,
        title=title,
        order_index=order,
        sha256=sha,
        extra={"cell_count": len(nb.cells)},
    )
    db.add(artifact)
    db.flush()

    meta = NOTEBOOK_META.get(path.name, {})
    notebook = Notebook(
        id=_id("nb", rel),
        artifact_id=art_id,
        filename=path.name,
        title=title,
        order_index=order,
        purpose=meta.get("purpose", ""),
        why_it_matters=meta.get("why", ""),
        expected_outcome=meta.get("outcome", ""),
    )
    db.add(notebook)
    db.flush()

    if lo_text:
        db.add(
            LearningObjective(
                id=_id("lo", rel),
                course_id="rad-llm",
                notebook_file=path.name,
                text=lo_text,
            )
        )

    current_heading = title
    for i, cell in enumerate(nb.cells):
        src = _cell_text(cell)
        out = _outputs(cell)
        h = _heading(src) if cell.cell_type == "markdown" else None
        if h:
            current_heading = h
        flags = inspect_code(src) if cell.cell_type == "code" else []
        evidence = EvidenceType.COURSE_SOURCE.value
        extra = {
            "safety_flags": flags,
            "has_stored_output": bool(out),
            "never_executed_by_academy": True,
        }
        if cell.cell_type == "code":
            extra["commands"] = _extract_commands(src)
            extra["models"] = _extract_models(src)
        span = SourceSpan(
            id=_id("sp", rel, i),
            artifact_id=art_id,
            source_type="notebook",
            file=path.name,
            cell_index=i,
            cell_type=cell.cell_type,
            heading=current_heading,
            body=_strip(src) if cell.cell_type == "markdown" else "",
            code=src if cell.cell_type == "code" else None,
            stored_output=out or None,
            execution_count=cell.get("execution_count"),
            evidence_type=evidence,
            embedding=hash_embedding((current_heading or "") + "\n" + src),
            extra=extra,
        )
        db.add(span)
        db.flush()
        db.add(
            NotebookCell(
                id=_id("cell", rel, i),
                notebook_id=notebook.id,
                span_id=span.id,
                cell_index=i,
                cell_type=cell.cell_type,
                source=src,
                stored_output=out or None,
                execution_count=cell.get("execution_count"),
                blocked_execution=True,
                safety_flags=flags,
                **_cell_commentary(path.name, i, cell.cell_type, src, flags),
            )
        )


def _extract_models(src: str) -> list[str]:
    found = re.findall(r"""(?:from_pretrained|pipeline)\(\s*['"]([^'"]+)['"]""", src)
    found += re.findall(r"""model\s*=\s*['"]([^'"]+)['"]""", src)
    found += re.findall(r"""['"]([A-Za-z0-9_.-]+/[A-Za-z0-9_.-]+)['"]""", src)
    return sorted(set(found))


def _extract_commands(src: str) -> list[str]:
    cmds = []
    for line in src.splitlines():
        s = line.strip()
        if s.startswith(("!", "%%")):
            cmds.append(s)
    return cmds


def _cell_commentary(
    filename: str, index: int, cell_type: str, src: str, flags: list[str]
) -> dict[str, str]:
    if cell_type != "code":
        return {
            "plain_english": "Expository markdown from the NVIDIA notebook.",
            "line_by_line": "",
            "why_exists": "Carries the course narrative, objectives, or wrap-up.",
            "what_should_happen": "The learner reads; nothing is executed.",
            "how_to_verify": "Match headings against the source notebook.",
            "common_failure": "Treating a diagram caption as an experimental measurement.",
            "try_modifying": "Rewrite the paragraph in school / engineer / research depth.",
        }
    danger = ", ".join(flags) if flags else "none flagged"
    return {
        "plain_english": (
            "Educational Python from the DLI notebook. The academy never runs it."
        ),
        "line_by_line": "\n".join(
            f"{n+1:03d}  {line}" for n, line in enumerate(src.splitlines()[:80])
        ),
        "why_exists": "Shows the HuggingFace / LangChain API the course wants you to internalize.",
        "what_should_happen": (
            "EXPECTED_RESULT only — this repo’s notebooks store no cell outputs."
        ),
        "how_to_verify": (
            "Re-run later in an isolated GPU sandbox; import telemetry as ACTUAL_RUN."
        ),
        "common_failure": (
            f"Assuming this cell already succeeded. Safety flags: {danger}."
        ),
        "try_modifying": "Predict shapes / tokens / memory before imagining a run.",
    }


def _ingest_pdf(db, path: Path) -> None:
    try:
        import fitz  # type: ignore
    except ImportError:
        return
    rel = str(path.relative_to(settings.course_materials_dir))
    art = SourceArtifact(
        id=_id("art", rel),
        course_id="rad-llm",
        source_type="pdf",
        filename=path.name,
        relpath=rel,
        title=path.stem,
        sha256=_sha(path),
    )
    db.add(art)
    doc = fitz.open(path)
    for i, page in enumerate(doc, start=1):
        text = page.get_text() or ""
        db.add(
            SourceSpan(
                id=_id("sp", rel, i),
                artifact_id=art.id,
                source_type="pdf",
                file=path.name,
                page=i,
                heading=f"Page {i}",
                body=text,
                embedding=hash_embedding(text),
                evidence_type=EvidenceType.COURSE_SOURCE.value,
            )
        )


def _ingest_pptx(db, path: Path) -> None:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return
    rel = str(path.relative_to(settings.course_materials_dir))
    art = SourceArtifact(
        id=_id("art", rel),
        course_id="rad-llm",
        source_type="pptx",
        filename=path.name,
        relpath=rel,
        title=path.stem,
        sha256=_sha(path),
    )
    db.add(art)
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
        body = "\n".join(chunks + ([notes] if notes else []))
        db.add(
            SourceSpan(
                id=_id("sp", rel, i),
                artifact_id=art.id,
                source_type="pptx",
                file=path.name,
                slide=i,
                heading=chunks[0][:200] if chunks else f"Slide {i}",
                body=body,
                embedding=hash_embedding(body),
                evidence_type=EvidenceType.COURSE_SOURCE.value,
            )
        )
